#!/usr/bin/env python3
"""Generate Kyanos investor presentation as PowerPoint with speaker notes."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_lines, notes_text):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle bullet points
        if line.startswith("• "):
            p.text = line[2:]
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(6)
        elif line.startswith("  - "):
            p.text = line[4:]
            p.level = 1
            p.font.size = Pt(16)
        elif line.startswith("**") and line.endswith("**"):
            p.text = line[2:-2]
            p.font.bold = True
            p.font.size = Pt(20)
            p.space_before = Pt(12)
        elif line == "":
            p.text = ""
            p.space_before = Pt(6)
        else:
            p.text = line
            p.font.size = Pt(18)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text

    return slide

def add_table_slide(prs, title, headers, rows, notes_text):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True

    # Table
    cols = len(headers)
    table_rows = len(rows) + 1
    table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * table_rows)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)

    # Speaker notes
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = notes_text

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # SLIDE 1: Title
    slide = add_title_slide(prs, "Kyanos", "Keeping Progressives Visible in the AI Era")
    notes = slide.notes_slide
    notes.notes_text_frame.text = """SPEAKER NOTES:

"Thanasi, thank you for your time. I reached out to Investing In US because you understand something most political operatives don't: the next communication revolution isn't coming—it's already here.

You and Reid founded Investing In US on the thesis that technology changes politics faster than politics adapts. That's exactly what's happening with AI.

Kyanos is the infrastructure play for progressive politics in the AI era. I want to show you why this belongs in your portfolio."

TRANSITION: "Let me start with a lesson we should have learned in 2024..." """

    # SLIDE 2: The Podcast Mistake
    add_content_slide(prs, "We Missed Podcasts. AI Is Next.", [
        "In 2020, progressives dismissed podcasts as entertainment for young men.",
        "",
        "Meanwhile, conservatives built an ecosystem that reached millions of persuadable voters every week.",
        "",
        "By the time we took it seriously, we were years behind.",
        "",
        "**The Pattern Repeats**",
        "• 2020-2024: Podcasts (Joe Rogan, Lex Fridman)",
        "• 2024-2028: AI Chatbots (Google AI Overviews, ChatGPT, WhatsApp AI)",
        "",
        "Progressives arrived late to podcasts. We have a chance to lead on AI."
    ], """SPEAKER NOTES:

"You watched this happen. Joe Rogan's endorsement of Trump reached 14 million people in a single episode. Progressives had no equivalent infrastructure.

The same pattern is emerging with AI. Right now, most campaigns treat AI as a curiosity—something the tech team might mention, not something the campaign manager prioritizes.

That's exactly how we treated podcasts in 2020."

TRANSITION: "But AI isn't just another podcast moment. The data shows it's more consequential..." """)

    # SLIDE 3: AI Information Revolution
    add_content_slide(prs, "By 2028, AI Will Answer Half of All Voter Queries", [
        "**Current State**",
        "• 58.5% of Google searches end without a click",
        "• 50% of Google queries show AI Overviews",
        "• 100M+ weekly ChatGPT users in US",
        "",
        "**What Matters Right Now**",
        "• Google AI Overviews - appears before any links",
        "• ChatGPT - 100M+ weekly US users actively researching",
        "• WhatsApp AI - critical channel for Hispanic communities",
        "",
        "**The Trend Is Clear**",
        "• Amazon Alexa+ launched Feb 2025 (Prime-bundled)",
        "• Apple partnering with Google for AI capabilities",
        "• Voice AI coming to kitchens, cars, living rooms"
    ], """SPEAKER NOTES:

"This isn't about some future scenario. 58% of Google searches already end without a click—voters get their answer from AI Overviews and never visit a campaign website.

By the time voters are making choices in 2028, nearly half of all political information queries will be answered inside AI environments, not on campaign websites.

The mission of the candidate website has transformed. It's no longer primarily a destination for voters—it's the source material that AI systems draw from when voters research."

TRANSITION: "And here's what makes this urgent: AI isn't just reaching voters. It's moving them—more effectively than any medium we've ever measured..." """)

    # SLIDE 4: The Research
    add_content_slide(prs, "AI Is 4x More Persuasive Than Television", [
        "**December 2025: Landmark Studies in Nature and Science**",
        "",
        "• Pro-Harris AI chatbot moved Trump supporters 3.9 points",
        "• 4x larger effect than TV advertising",
        "• Canadian and Polish elections: ~10 percentage points shift",
        "• AI optimized for persuasion: up to 25 percentage points",
        "",
        '"This is a shockingly large effect... especially in the context of presidential politics." — David Rand, MIT',
        "",
        "**Meanwhile, Traditional Contact Has Near-Zero Effect**",
        "",
        "Kalla & Broockman (2018): 49 field experiments showed campaign contact and advertising have approximately zero effect on candidate choice."
    ], """SPEAKER NOTES:

"Thanasi, this is peer-reviewed research from the world's most prestigious scientific journals, published just last month.

AI chatbots are the most persuasive political communication channel ever measured in controlled experiments. Four times more effective than television advertising.

And the beauty of Internet Presence Optimization: it costs a fraction of TV. We're talking $1-10 per voter influenced, compared to $100-600 for broadcast television.

This isn't just about keeping pace. It's about finding a lever that actually works when traditional contact doesn't."

TRANSITION: "But there's another dimension to this that I know you've been thinking about..." """)

    # SLIDE 5: Platform Risk
    add_content_slide(prs, "AI Platforms Face Structural Pressure", [
        "**The Concentration Problem**",
        "• Microsoft partners exclusively with OpenAI",
        "• Amazon and Google invested billions in Anthropic",
        "• Meta distributes Llama at platform scale",
        "",
        "**The Leverage Points**",
        "• Government contracts (DoD JWCC: $9B, IC C2E: tens of billions)",
        "• Regulatory exposure (antitrust, Section 230, mergers)",
        "• Political pressure from administration demanding loyalty",
        "",
        "**The Risk**",
        "This isn't about tech leaders' intentions—many are personally progressive.",
        "",
        "The issue is structural: responsible corporate governance requires managing political risk. In the current environment, that risk runs in one direction."
    ], """SPEAKER NOTES:

"You know these companies better than I do. You know the pressures they face.

I'm not suggesting bad faith. I'm suggesting that the incentive structures are misaligned. When an administration is willing to punish perceived disloyalty—and when companies have billions in government contracts at stake—caution has a political direction.

That's why external monitoring matters. Campaigns and the progressive movement can't rely on platforms to self-police. We need independent visibility into how AI systems represent political topics—and the ability to sound the alarm if something changes.

Kyanos provides that watchdog function."

TRANSITION: "So what exactly are we building?" """)

    # SLIDE 6: The Solution
    add_content_slide(prs, "Kyanos: Monitor, Diagnose, Fix", [
        "A family of AI-driven services that help progressive campaigns manage their digital presence.",
        "",
        "**What We Do**",
        "• Monitor how AI systems represent progressive candidates",
        "• Trace errors to their root causes",
        "• Diagnose what's actually fixable",
        "• Prescribe specific actions campaigns can implement immediately",
        "",
        "**What We Don't Do**",
        "• Tell campaigns what they 'should have' done",
        "• Generate reports that sit unread",
        "• Recommend things outside their control",
        "",
        "**The Watchdog Role**",
        "If models shift against progressives, we know first and sound the alarm."
    ], """SPEAKER NOTES:

"The name Kyanos comes from the ancient Greek word for blue—the root of 'cyan.' We chose it because we're building something foundational: the infrastructure that keeps progressives visible in the AI era.

Our methodology is simple: Observe what AI says, Trace where it got that information, Diagnose the root cause, and Prescribe specific fixes.

We don't just hand over a report. We monitor execution, coach the team through implementation, and confirm when AI systems update their responses."

TRANSITION: "Let me show you how this works in practice..." """)

    # SLIDE 7: How It Works
    add_content_slide(prs, "The Causality Chain Methodology", [
        "**Phase 1: Monitor**",
        "• Query AI surfaces (Google AI Overviews, ChatGPT, Meta AI)",
        "• Test questions voters actually ask about the candidate",
        "• Watchtower detects geographic variation",
        "",
        "**Phase 2: Prioritize**",
        "• Trace errors to root causes (Wikipedia? News? Campaign site?)",
        "• Rank by: Voter Reach × Error Severity × Fixability",
        "• Focus only on things the campaign can change",
        "",
        "**Phase 3: Act**",
        "• Step-by-step instructions for each issue",
        "• Tailored to campaign's CMS (WordPress, Squarespace, NationBuilder)",
        "• Monitor execution and coach through completion"
    ], """SPEAKER NOTES:

"Here's a concrete example. ChatGPT says a candidate supports a policy they actually oppose. We don't just flag the error—we trace it.

Maybe it's a Wikipedia entry that needs correcting. Maybe it's a misquoted news article that's been indexed. Maybe the campaign's issues page is missing or poorly structured.

Then we tell the campaign exactly what needs to happen. The comms director gets the messaging guidance. The web developer gets the technical specs. We provide templates, code snippets, everything they need to execute.

And we don't walk away. We track implementation and confirm when AI systems update."

TRANSITION: "Beyond Panopticon, we have a complete product suite..." """)

    # SLIDE 8: Product Suite
    slide = add_table_slide(prs, "Four Products, One Platform",
        ["Product", "Function", "Status"],
        [
            ["Panopticon", "Core AI monitoring and remediation", "MVP+"],
            ["GoodInk", "Positive press discovery for campaigns to feature", "MVP"],
            ["Adwatch", "Political ad monitoring across Google and Meta", "MVP"],
            ["TripWire", "Opponent website change detection", "MVP"],
        ],
        """SPEAKER NOTES:

"Panopticon is the flagship—it's where campaigns see what AI says about them and get fixes.

GoodInk solves a different problem: campaigns often miss positive press coverage they could be featuring. We surface that automatically.

Adwatch monitors the political advertising landscape—what opponents are running, where, and how it's targeting voters.

TripWire watches opponent websites for changes. When they update their messaging or positions, campaigns know immediately.

All of this runs on Bentham, our extraction infrastructure. It's named after Jeremy Bentham, the philosopher who designed the original Panopticon."

ADDITIONAL:
• Watchtower: Distributed observation network for geographic monitoring
• Bentham Infrastructure: The extraction engine powering all products

TRANSITION: "Now let me explain how we've structured this to be permanent and accountable..." """)

    # SLIDE 9: How We Operate
    add_content_slide(prs, "How We Operate", [
        "**Trust-Owned. Mission-Locked. Sustainable.**",
        "",
        "**Structure**",
        "• Trust entity owns 100% of Kyanos LLC",
        "• Board of trustees (3-5 progressive leaders) governs trust",
        "• Trust charter locks in progressive-only service",
        "• Prevents hostile acquisition—can never be sold to opponents",
        "",
        "**The Catalist Precedent**",
        "• Founded 2006 with $1M from Soros, $4M+ from Tides",
        "• Trust-owned LLC structure",
        "• Today serves virtually every major Democratic campaign",
        "",
        "**The model is proven. Foundation funding creates permanent progressive infrastructure.**"
    ], """SPEAKER NOTES:

"I know you've seen a lot of pitch decks. Most promise equity and exits. This is different.

We're not optimizing for an exit. We're building infrastructure that the progressive movement will rely on for decades. The trust structure ensures that.

Catalist proved this model works. They took foundation funding in 2006, built permanent infrastructure, and now serve virtually every major Democratic campaign. We're following the same playbook.

Your investment doesn't come with equity—there is no equity. What you get is a seat at the table for governance, and the knowledge that you helped build something permanent."

TRANSITION: "Here's the market we're going after..." """)

    # SLIDE 10: Market & Go-to-Market
    add_content_slide(prs, "The Path to 466 Clients", [
        "**2026: The Free Year (Foundation-Funded)**",
        "• 15 Statewide campaigns",
        "• 50 Congressional campaigns",
        "• 110 Local campaigns",
        "• 15 Officeholders + 10 Organizations",
        "• Total: 200 campaigns served free",
        "",
        "**Purpose:** Build reference base, case studies, testimonials",
        "",
        "**2027-2028: Paid Launch**",
        "• 2027: 150 clients, $1.2M revenue, Breakeven",
        "• 2028: 466 clients, $2.8M revenue, +$850K net",
        "",
        "**Go-to-Market:** Consultants are primary channel for larger races"
    ], """SPEAKER NOTES:

"The 2026 free year is strategic. We're investing $330K of the raise to serve 200 campaigns at no cost. That builds our reference base and generates the case studies we need to sell in 2027.

The go-to-market is deliberately low-friction. Small campaigns get self-service—enter your info, system configures automatically. Larger campaigns come through consultants who already have trusted relationships.

By 2028, we project 466 clients and $2.8M in revenue. We're cash-flow positive by late 2027."

TRANSITION: "You're going to ask about competition..." """)

    # SLIDE 11: Competitive Position
    add_content_slide(prs, "First Mover in Progressive AI Monitoring", [
        "**The Landscape**",
        "• Generic SEO tools (SEMrush, Ahrefs): Not political, not AI-focused",
        "• AI monitoring: Mostly enterprise brand monitoring",
        "• Political tech (NGP VAN, Catalist): Voter data, not AI presence",
        "• Progressive AI: Kyanos is the only player focused here",
        "",
        "**What Makes Us Different**",
        "• Political context—we understand campaigns",
        "• AI-native—optimized for AI answer generation",
        "• Causality focus—trace problems to root causes",
        "• Actionable fixes—solutions, not just reports",
        "• Progressive-only—trust structure locks in mission"
    ], """SPEAKER NOTES:

"You'll notice there's no direct competitor. That's because the category doesn't exist yet.

SEO tools aren't built for AI. AI monitoring tools aren't built for politics. Political tech companies are focused on voter data, not AI presence.

We're the first mover in progressive AI monitoring, and our trust structure means we can't be acquired by someone who would use this against us.

The network effects matter here. Every campaign we serve makes our intelligence better. Every fix we implement teaches us more about how AI systems source information."

TRANSITION: "Let me introduce the team..." """)

    # SLIDE 12: Team & Governance
    add_content_slide(prs, "Built for Permanence", [
        "**Founder**",
        "Ed Forman — 35 years Silicon Valley experience. Early employee at eleven startups. MBA from Stanford, later taught entrepreneurship there.",
        "",
        "Role: Architect and advisor. Not taking an operational role.",
        "",
        "**Launch Team (Q2 2026)**",
        "• CEO: Fundraising, foundation relationships, strategy",
        "• Technical Lead: Bentham infrastructure, product development",
        "• Research Director: Methodology, AI Bias Report, quality",
        "• Operations Manager: Study execution, client delivery",
        "• Marketing & Partnerships: Consultant relationships, sales"
    ], """SPEAKER NOTES:

"I'm the architect, not the operator. My role is to assemble the team, frame the problem, and design a structure that outlasts any individual.

We'll hire a CEO with deep relationships in the progressive funding community. The technical lead will build Bentham and the product suite. Research director ensures methodology rigor.

The board of trustees governs the trust and ensures mission fidelity. Major funders get board seats—that's part of the value proposition for foundation investors."

TRANSITION: "Here's what we're asking for..." """)

    # SLIDE 13: The Ask
    add_content_slide(prs, "$1,000,000 to Build Progressive AI Infrastructure", [
        "**The Ask**",
        "• Target: $1,000,000",
        "• Minimum: $750,000",
        "• Reach: $1,250,000",
        "• Timeline: Commitments by March 15, 2026",
        "",
        "**Use of Funds**",
        "• Personnel (5 FTEs): $650,000 (65%)",
        "• Contractors: $150,000 (15%)",
        "• Infrastructure: $48,000 (5%)",
        "• Marketing/Events: $40,000 (4%)",
        "• Legal/Operations/Reserve: $112,000 (11%)",
        "",
        "**Self-sustaining by 2028 with +$850K annual surplus**"
    ], """SPEAKER NOTES:

"We're asking for $1 million. That funds a full year of building and a free access program for 200 campaigns in 2026.

Why can $1 million do what would have taken $5 million a few years ago? AI-assisted development has fundamentally changed the economics. We practice what we preach—Kyanos is built with the same AI capabilities we help campaigns leverage.

The timeline is tight by design. We need commitments by March 15 to be operational by June and serving campaigns before the 2026 cycle heats up."

TRANSITION: "Let me tell you why I'm in this room specifically..." """)

    # SLIDE 14: Why Investing In US
    add_content_slide(prs, "This Is What You Were Built For", [
        '**Your Core Thesis—Applied to AI**',
        '"Technology is changing politics faster than politics is adapting to technology."',
        "— Investing In US founding principle",
        "",
        "**That's exactly what's happening with AI and voter information.**",
        "",
        "**You've Already Invested in This Space**",
        "• Vote.org: Voter registration → We ensure voters find accurate info",
        "• News for Democracy: Media → We ensure AI surfaces accurate news",
        "• Run for Something: Candidates → We help them compete in AI",
        "",
        "**Kyanos completes your portfolio. It's the AI layer on top of everything else.**",
        "",
        "Help us keep progressives visible in the AI era."
    ], """SPEAKER NOTES:

"Thanasi, Investing In US was founded on a simple observation: technology is changing politics faster than politics is adapting.

That's exactly what we're seeing with AI. Political operatives are treating chatbots like a curiosity while conservatives are already optimizing for them.

You've already invested in voter registration, candidate pipelines, and media ecosystems. Kyanos is the missing layer—ensuring that when voters use AI to research candidates, they find accurate information.

And here's what I know about Investing In US: you're not looking for financial returns. You're looking for democratic returns. You want to build infrastructure that outlasts any single campaign or election cycle.

That's what Kyanos is. Trust-owned. Mission-locked. Built to serve progressives permanently.

The window is now. By 2028, the patterns will be set. The question is whether progressives have infrastructure in place—or whether we're playing catch-up on AI like we did with podcasts.

Here's my ask: I'd like you to bring this to Reid and Dmitri. If this resonates with you—if you see how it completes your portfolio—I want to give you everything you need to make the case internally.

What questions do you have? What would they need to see?" """)

    # Save
    prs.save('/Users/edf/KyanosTech-Plan/docs/presentations/kyanos-investor-pitch.pptx')
    print("Created: kyanos-investor-pitch.pptx")

if __name__ == "__main__":
    main()
