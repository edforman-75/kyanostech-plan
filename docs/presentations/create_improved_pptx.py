#!/usr/bin/env python3
"""Generate improved Kyanos presentation with Duarte suggestions."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Color scheme
BLUE_DARK = RgbColor(0x1a, 0x36, 0x5d)
BLUE_MID = RgbColor(0x00, 0x7b, 0xc0)
BLUE_LIGHT = RgbColor(0x4a, 0xb8, 0xe8)
ORANGE = RgbColor(0xff, 0x6b, 0x35)
GREEN = RgbColor(0x2e, 0xcc, 0x71)
GRAY = RgbColor(0x6c, 0x75, 0x7d)
WHITE = RgbColor(0xff, 0xff, 0xff)

def set_shape_color(shape, fill_color, line_color=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()

def add_big_stat(slide, number, label, x, y, width=3, color=BLUE_MID):
    num_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.1), Inches(width), Inches(0.8))
    tf = label_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: TITLE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.3))
    set_shape_color(bar, BLUE_MID)

    # Logo
    try:
        slide.shapes.add_picture('/Users/edf/KyanosTech-Plan/docs/presentations/logos/kyanos-main.png',
                                  Inches(5.7), Inches(1.2), height=Inches(2.2))
    except:
        logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(1.5), Inches(1.5), Inches(1.5))
        set_shape_color(logo, BLUE_MID)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    # Updated tagline per Duarte
    tag = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.6))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "When Voters Ask, What Will AI Say?"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BLUE_MID
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(12.333), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Monitoring for Progressive Campaigns"
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    info = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "Ed Forman, Founder  •  For Thanasi Dilos, Investing In US  •  January 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Thanasi, thank you for your time. When voters want to learn about a candidate, they don't visit campaign websites anymore—they ask AI.

Kyanos ensures AI gives them the right answer."

TRANSITION: "Let me start with a lesson we should have learned in 2024..." """)

    # ========== SLIDE 2: PODCAST MISTAKE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "We Missed Podcasts. AI Is Next."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Left - THEN (red)
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
    set_shape_color(box1, RgbColor(0xff, 0xeb, 0xee))
    box1.adjustments[0] = 0.05

    header1 = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(5.5), Inches(0.6))
    tf = header1.text_frame
    p = tf.paragraphs[0]
    p.text = "2020-2024: PODCASTS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    stat1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(5.5), Inches(1))
    tf = stat1.text_frame
    p = tf.paragraphs[0]
    p.text = "14M"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    label1 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.8))
    tf = label1.text_frame
    p = tf.paragraphs[0]
    p.text = "viewers saw Rogan's\nTrump endorsement"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    result1 = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(5.5), Inches(1))
    tf = result1.text_frame
    p = tf.paragraphs[0]
    p.text = "Progressives arrived late"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.4), Inches(3.8), Inches(0.8), Inches(0.5))
    set_shape_color(arrow, BLUE_MID)

    # Right - NOW (green)
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.5), Inches(4.5))
    set_shape_color(box2, RgbColor(0xe8, 0xf5, 0xe9))
    box2.adjustments[0] = 0.05

    header2 = slide.shapes.add_textbox(Inches(7.2), Inches(2), Inches(5.5), Inches(0.6))
    tf = header2.text_frame
    p = tf.paragraphs[0]
    p.text = "2024-2028: AI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # AI platforms
    platforms = slide.shapes.add_textbox(Inches(7.2), Inches(2.7), Inches(5.5), Inches(1.5))
    tf = platforms.text_frame
    p = tf.paragraphs[0]
    p.text = "Google AI • ChatGPT • Meta AI"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    insight = slide.shapes.add_textbox(Inches(7.2), Inches(3.5), Inches(5.5), Inches(1.2))
    tf = insight.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Three companies control\nmost AI information access"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    result2 = slide.shapes.add_textbox(Inches(7.2), Inches(5), Inches(5.5), Inches(1))
    tf = result2.text_frame
    p = tf.paragraphs[0]
    p.text = "We can lead this time"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Bottom quote
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.8))
    set_shape_color(quote, BLUE_DARK)
    quote.adjustments[0] = 0.2
    quote_txt = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.5))
    tf = quote_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "History doesn't repeat, but it rhymes — we have one chance to get ahead."
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"You watched this happen. Joe Rogan's endorsement reached 14 million people. Progressives had no equivalent infrastructure.

The same pattern is emerging with AI. Three companies—Google, OpenAI, Meta—control most AI information access.

History doesn't repeat, but it rhymes. We have one chance to get ahead."

TRANSITION: "But AI isn't just another podcast moment. The data shows it's more consequential..." """)

    # ========== SLIDE 3: AI STATS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Controls the Answers"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Big stats
    add_big_stat(slide, "58%", "of searches end\nwithout a click", 0.5, 1.3, 3)
    add_big_stat(slide, "50%", "show AI\nOverviews", 3.7, 1.3, 3)
    add_big_stat(slide, "100M+", "weekly ChatGPT\nusers in US", 6.9, 1.3, 3)
    add_big_stat(slide, "45-60%", "voter queries\nby AI in 2028", 10.1, 1.3, 3, ORANGE)

    # Key insight box
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4), Inches(11.8), Inches(1.2))
    set_shape_color(insight_box, RgbColor(0xe3, 0xf2, 0xfd))
    insight_box.adjustments[0] = 0.1

    insight = slide.shapes.add_textbox(Inches(0.8), Inches(4.25), Inches(11.8), Inches(0.8))
    tf = insight.text_frame
    p = tf.paragraphs[0]
    p.text = "When AI controls the answers, controlling what AI learns becomes critical."
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    # Platform boxes
    box_y = 5.5
    platforms_data = [
        ("Google AI", "Appears before all links", RgbColor(0xea, 0x43, 0x35)),
        ("ChatGPT", "100M+ active researchers", RgbColor(0x10, 0xa3, 0x7f)),
        ("Meta AI / WhatsApp", "Critical for Hispanic voters", RgbColor(0x00, 0x84, 0xff)),
    ]

    for i, (name, desc, color) in enumerate(platforms_data):
        x = 0.8 + i * 4.2
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(box_y), Inches(3.8), Inches(1.5))
        set_shape_color(box, color)
        box.adjustments[0] = 0.15

        name_txt = slide.shapes.add_textbox(Inches(x), Inches(box_y + 0.25), Inches(3.8), Inches(0.5))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        desc_txt = slide.shapes.add_textbox(Inches(x), Inches(box_y + 0.8), Inches(3.8), Inches(0.5))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"58% of Google searches end without a click—voters get their answer from AI and never visit a campaign website.

When AI controls the answers, controlling what AI learns becomes critical. That's what Kyanos does.

Three platforms matter now: Google AI Overviews, ChatGPT, and Meta AI—especially WhatsApp for Hispanic communities."

TRANSITION: "And here's what makes this urgent: AI isn't just reaching voters. It's moving them..." """)

    # ========== SLIDE 4: PERSUASION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "4x More Persuasive at 1/60th the Cost"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    source = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(0.4), Inches(3.3), Inches(0.5))
    set_shape_color(source, RgbColor(0xe3, 0xf2, 0xfd))
    source.adjustments[0] = 0.3
    src_text = slide.shapes.add_textbox(Inches(9.5), Inches(0.45), Inches(3.3), Inches(0.4))
    tf = src_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Nature & Science, Dec 2025"
    p.font.size = Pt(12)
    p.font.color.rgb = BLUE_MID
    p.alignment = PP_ALIGN.CENTER

    # TV bar
    tv_label = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(2), Inches(0.4))
    tf = tv_label.text_frame
    p = tf.paragraphs[0]
    p.text = "TV Advertising"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    tv_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(2.5), Inches(0.8))
    set_shape_color(tv_bar, GRAY)
    tv_bar.adjustments[0] = 0.2

    tv_val = slide.shapes.add_textbox(Inches(1), Inches(2.1), Inches(2.5), Inches(0.6))
    tf = tv_val.text_frame
    p = tf.paragraphs[0]
    p.text = "1x"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # AI bar
    ai_label = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(2), Inches(0.4))
    tf = ai_label.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Chatbots"
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE_MID

    ai_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3.7), Inches(10), Inches(0.8))
    set_shape_color(ai_bar, BLUE_MID)
    ai_bar.adjustments[0] = 0.05

    ai_val = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(10), Inches(0.6))
    tf = ai_val.text_frame
    p = tf.paragraphs[0]
    p.text = "4x MORE PERSUASIVE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Cost comparison
    cost_y = 5
    tv_cost = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(cost_y), Inches(5), Inches(1.8))
    set_shape_color(tv_cost, RgbColor(0xf5, 0xf5, 0xf5))
    tv_cost.adjustments[0] = 0.1

    tv_price = slide.shapes.add_textbox(Inches(1), Inches(cost_y + 0.2), Inches(5), Inches(0.8))
    tf = tv_price.text_frame
    p = tf.paragraphs[0]
    p.text = "$100-600"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    tv_desc = slide.shapes.add_textbox(Inches(1), Inches(cost_y + 1.1), Inches(5), Inches(0.4))
    tf = tv_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "per voter influenced (TV)"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    ai_cost = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(cost_y), Inches(5), Inches(1.8))
    set_shape_color(ai_cost, RgbColor(0xe8, 0xf5, 0xe9))
    ai_cost.adjustments[0] = 0.1

    ai_price = slide.shapes.add_textbox(Inches(7), Inches(cost_y + 0.2), Inches(5), Inches(0.8))
    tf = ai_price.text_frame
    p = tf.paragraphs[0]
    p.text = "$1-10"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    ai_desc = slide.shapes.add_textbox(Inches(7), Inches(cost_y + 1.1), Inches(5), Inches(0.4))
    tf = ai_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "per voter influenced (AI)"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Quote
    quote = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(8), Inches(0.4))
    tf = quote.text_frame
    p = tf.paragraphs[0]
    p.text = '"A shockingly large effect" — David Rand, MIT'
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = GRAY

    add_notes(slide, """SPEAKER NOTES:

"MIT research published in Nature confirms AI chatbots are 4x more persuasive than TV advertising. David Rand called it 'a shockingly large effect.'

And the cost? $1-10 per voter influenced vs $100-600 for TV. That's 4x the persuasion at 1/60th the cost.

This makes AI the dominant battleground for 2026 and beyond."

TRANSITION: "But there's another dimension to this..." """)

    # ========== SLIDE 5: PLATFORM RISK ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Three Pressure Points on AI Platforms"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Three risk boxes
    risks = [
        ("$9B+", "Government\nContracts", "DoD, Intelligence Community"),
        ("230", "Section 230\nExposure", "Liability, antitrust, mergers"),
        ("⚠", "Political\nPressure", "Administration demands loyalty"),
    ]

    for i, (icon, title_text, desc) in enumerate(risks):
        x = 1 + i * 4
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(3.5), Inches(3))
        set_shape_color(box, RgbColor(0xff, 0xf3, 0xe0))
        box.adjustments[0] = 0.1

        ico = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(3.5), Inches(0.8))
        tf = ico.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER

        ttl = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(3.5), Inches(0.8))
        tf = ttl.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        dsc = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(3.5), Inches(0.5))
        tf = dsc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Key message
    msg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5), Inches(11.8), Inches(2))
    set_shape_color(msg_box, BLUE_DARK)
    msg_box.adjustments[0] = 0.1

    msg = slide.shapes.add_textbox(Inches(0.8), Inches(5.3), Inches(11.8), Inches(1.5))
    tf = msg.text_frame
    p = tf.paragraphs[0]
    p.text = "Without independent monitoring,"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "progressive voices become invisible by design."
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Three pressure points: $9B+ in government contracts. Section 230 liability exposure. And an administration that demands loyalty.

I'm not suggesting bad faith. But responsible corporate governance requires managing political risk. In this environment, that risk runs one direction.

Without independent monitoring, progressive voices become invisible by design. That's why Kyanos exists."

TRANSITION: "So what exactly are we building?" """)

    # ========== SLIDE 6: SOLUTION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kyanos: Monitor → Trace → Diagnose → Fix"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # 4-step process
    steps = [
        ("🔍", "MONITOR", "See what AI\nsays about you"),
        ("🔗", "TRACE", "Find where it\ngot that info"),
        ("📊", "DIAGNOSE", "Rank issues\nby impact"),
        ("✓", "FIX", "CMS-specific\naction steps"),
    ]

    for i, (icon, title_text, desc) in enumerate(steps):
        x = 0.8 + i * 3.2

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.8), Inches(1.5), Inches(1.4), Inches(1.4))
        set_shape_color(circle, BLUE_MID)

        icon_txt = slide.shapes.add_textbox(Inches(x + 0.8), Inches(1.8), Inches(1.4), Inches(0.8))
        tf = icon_txt.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(40)
        p.alignment = PP_ALIGN.CENTER

        step_title = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(3), Inches(0.5))
        tf = step_title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        step_desc = slide.shapes.add_textbox(Inches(x), Inches(3.6), Inches(3), Inches(0.8))
        tf = step_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.8), Inches(2), Inches(0.5), Inches(0.4))
            set_shape_color(arrow, RgbColor(0xcc, 0xcc, 0xcc))

    # Differentiators
    diff_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.3))
    set_shape_color(diff_box, RgbColor(0xf5, 0xf5, 0xf5))
    diff_box.adjustments[0] = 0.08

    diff_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.65), Inches(5.5), Inches(0.4))
    tf = diff_title.text_frame
    p = tf.paragraphs[0]
    p.text = "What Makes Us Different"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    diffs = ["Campaign-focused, not generic SEO", "AI-native, not search rankings", "Root causes, not just reports", "Ready-to-execute fixes"]
    for i, diff in enumerate(diffs):
        txt = slide.shapes.add_textbox(Inches(1.2), Inches(5.1 + i * 0.4), Inches(5), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = "✓ " + diff
        p.font.size = Pt(14)
        p.font.color.rgb = GREEN

    # Watchdog box
    watch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.3))
    set_shape_color(watch_box, BLUE_DARK)
    watch_box.adjustments[0] = 0.08

    watch_title = slide.shapes.add_textbox(Inches(6.8), Inches(4.65), Inches(5.8), Inches(0.4))
    tf = watch_title.text_frame
    p = tf.paragraphs[0]
    p.text = "🔔 Early Warning System"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    watch_desc = slide.shapes.add_textbox(Inches(6.8), Inches(5.2), Inches(5.8), Inches(1.2))
    tf = watch_desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "If AI models shift against progressives, we detect it first and sound the alarm."
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Four steps: Monitor what AI says. Trace where it got that information. Diagnose root causes and rank by impact. Then provide specific fixes campaigns can execute immediately.

We're not generic SEO—we're campaign-focused. We don't just generate reports—we provide ready-to-execute action items.

And if AI models shift against progressives, we detect it first and sound the alarm."

TRANSITION: "Let me show you how this works in practice..." """)

    # ========== SLIDE 7: PRACTICAL EXAMPLE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "From AI Confusion to Clear Messaging in 48 Hours"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Query box
    q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3), Inches(5.5), Inches(1))
    set_shape_color(q_box, BLUE_LIGHT)
    q_box.adjustments[0] = 0.15
    q_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6))
    tf = q_txt.text_frame
    p = tf.paragraphs[0]
    p.text = '💬 "What does [Candidate] think about abortion?"'
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    a1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.3), Inches(2.4), Inches(0.5), Inches(0.4))
    set_shape_color(a1, GRAY)

    # Problem
    r_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.9), Inches(5.5), Inches(1.2))
    set_shape_color(r_box, RgbColor(0xff, 0xcc, 0xbc))
    r_box.adjustments[0] = 0.1
    r_txt = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(5.5), Inches(1))
    tf = r_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ AI Response:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p = tf.add_paragraph()
    p.text = '"The candidate has not stated a clear position..."'
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY

    # Arrow
    a2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.3), Inches(4.2), Inches(0.5), Inches(0.4))
    set_shape_color(a2, GRAY)

    # Root cause
    c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.7), Inches(5.5), Inches(1.2))
    set_shape_color(c_box, RgbColor(0xff, 0xf9, 0xc4))
    c_box.adjustments[0] = 0.1
    c_txt = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(5.5), Inches(1))
    tf = c_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "🔍 Root Cause:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p = tf.add_paragraph()
    p.text = "Issues page missing structured data markup"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # Right side - Solution
    fix_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5))
    set_shape_color(fix_box, RgbColor(0xe8, 0xf5, 0xe9))
    fix_box.adjustments[0] = 0.05

    fix_title = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    tf = fix_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ THE FIX"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    fix_steps = [
        "1. Add structured FAQ to issues page",
        "2. Include clear position statement",
        "3. Add schema.org markup",
        "4. Submit to Google Search Console",
    ]

    for i, step in enumerate(fix_steps):
        step_txt = slide.shapes.add_textbox(Inches(7.1), Inches(2.2 + i * 0.7), Inches(5.2), Inches(0.5))
        tf = step_txt.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE_DARK

    # Result
    result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(5.2), Inches(5.2), Inches(1.2))
    set_shape_color(result_box, GREEN)
    result_box.adjustments[0] = 0.15

    result = slide.shapes.add_textbox(Inches(7.1), Inches(5.5), Inches(5.2), Inches(0.6))
    tf = result.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ AI now shows correct position"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Here's a concrete example. Voter asks AI about a candidate's abortion position. AI says 'no clear position stated'—but the candidate has a clear position.

We trace it: the issues page exists but lacks structured data AI can parse.

The fix: Add structured FAQ, clear position statement, schema.org markup, submit to Google.

Result: From confusion to clarity in 48 hours."

TRANSITION: "Beyond Panopticon, we have a complete product suite..." """)

    # ========== SLIDE 8: PRODUCTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Kyanos Product Suite"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    products = [
        ("PANOPTICON", "MVP+", "AI Monitoring & Remediation", "Core product • Launching Q2 2026", BLUE_MID, '/Users/edf/KyanosTech-Plan/docs/presentations/logos/panopticon.png'),
        ("GOODINK", "MVP", "Positive Press Discovery", "Surface coverage worth featuring", RgbColor(0x00, 0x96, 0x88), '/Users/edf/KyanosTech-Plan/docs/presentations/logos/goodink.png'),
        ("ADWATCH", "MVP", "Ad Monitoring", "Track opponent ads across platforms", RgbColor(0x7b, 0x1f, 0xa2), '/Users/edf/KyanosTech-Plan/docs/presentations/logos/adwatch.png'),
        ("TRIPWIRE", "MVP", "Opponent Monitoring", "Know when competitors shift messaging", RgbColor(0xf5, 0x7c, 0x00), '/Users/edf/KyanosTech-Plan/docs/presentations/logos/tripwire.png'),
    ]

    for i, (name, status, subtitle, desc, color, logo_path) in enumerate(products):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.3
        y = 1.2 + row * 2.7

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.4))
        set_shape_color(card, RgbColor(0xfa, 0xfa, 0xfa), color)
        card.adjustments[0] = 0.08
        card.line.width = Pt(3)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.15), Inches(2.4))
        set_shape_color(bar, color)

        # Try to add logo
        try:
            slide.shapes.add_picture(logo_path, Inches(x + 0.3), Inches(y + 0.3), height=Inches(0.6))
        except:
            pass

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 4.5), Inches(y + 0.2), Inches(1), Inches(0.4))
        badge_color = GREEN if status == "MVP+" else RgbColor(0x90, 0xa4, 0xae)
        set_shape_color(badge, badge_color)
        badge.adjustments[0] = 0.4
        badge_txt = slide.shapes.add_textbox(Inches(x + 4.5), Inches(y + 0.25), Inches(1), Inches(0.3))
        tf = badge_txt.text_frame
        p = tf.paragraphs[0]
        p.text = status
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        name_txt = slide.shapes.add_textbox(Inches(x + 1), Inches(y + 0.3), Inches(3.3), Inches(0.5))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color

        sub_txt = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 1), Inches(5), Inches(0.4))
        tf = sub_txt.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK

        desc_txt = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 1.5), Inches(5), Inches(0.6))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    # Watchtower callout
    watch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.9))
    set_shape_color(watch, RgbColor(0x00, 0x60, 0x64))
    watch.adjustments[0] = 0.15

    watch_txt = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.8), Inches(0.5))
    tf = watch_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "🗼 WATCHTOWER: Geographic monitoring detects regional AI variation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Panopticon is the flagship—MVP+, launching Q2 2026. It's where campaigns see what AI says and get fixes.

GoodInk surfaces positive press coverage. AdWatch monitors opponent advertising. TripWire detects when competitors change messaging.

Watchtower is our geographic monitoring network—it detects regional AI variation, critical for statewide and national campaigns."

TRANSITION: "Now let me explain how we've structured this to be permanent and accountable..." """)

    # ========== SLIDE 9: COMPETITIVE POSITION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "First Mover in Progressive AI Monitoring"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    headers = ["", "Political\nContext", "AI Bias\nDetection", "Root\nCauses", "Actionable\nFixes", "Progressive\nOnly"]
    for i, h in enumerate(headers):
        x = 0.8 + i * 2
        w = 1.8 if i > 0 else 2.5
        txt = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(w), Inches(0.8))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

    competitors = [
        ("SEMrush / Ahrefs", ["❌", "❌", "❌", "❌", "❌"]),
        ("AI Monitoring Tools", ["❌", "❌", "❌", "❌", "❌"]),
        ("NGP VAN / Catalist", ["✓", "❌", "❌", "❌", "✓"]),
        ("KYANOS", ["✓", "✓", "✓", "✓", "✓"]),
    ]

    for row, (name, checks) in enumerate(competitors):
        y = 2 + row * 1
        is_kyanos = name == "KYANOS"

        if is_kyanos:
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y - 0.1), Inches(12.2), Inches(0.9))
            set_shape_color(bg, RgbColor(0xe3, 0xf2, 0xfd))
            bg.adjustments[0] = 0.15

        name_txt = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.1), Inches(2.5), Inches(0.6))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = is_kyanos
        p.font.color.rgb = BLUE_MID if is_kyanos else GRAY

        for i, check in enumerate(checks):
            x = 3.3 + i * 2
            check_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.05), Inches(1.8), Inches(0.6))
            tf = check_txt.text_frame
            p = tf.paragraphs[0]
            p.text = check
            p.font.size = Pt(24)
            p.font.color.rgb = GREEN if check == "✓" else RgbColor(0xcc, 0xcc, 0xcc)
            p.alignment = PP_ALIGN.CENTER

    # Network effects
    net_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.4))
    set_shape_color(net_box, BLUE_DARK)
    net_box.adjustments[0] = 0.1

    net_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.4))
    tf = net_title.text_frame
    p = tf.paragraphs[0]
    p.text = "NETWORK EFFECTS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xbb, 0xde, 0xfb)
    p.alignment = PP_ALIGN.CENTER

    net_desc = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.6))
    tf = net_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "More campaigns → Better intelligence → Stronger platform accountability"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"There's no direct competitor because the category doesn't exist yet.

SEO tools aren't built for AI or politics. AI monitoring tools don't understand campaigns. Political tech focuses on voter data, not AI presence.

We're first mover, and our trust structure means we can't be acquired by opponents.

Network effects matter: more campaigns means better intelligence means stronger accountability."

TRANSITION: "Let me explain our governance structure..." """)

    # ========== SLIDE 10: GOVERNANCE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Trust-Owned. Mission-Locked. Permanent."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Structure diagram
    trust_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(1.2), Inches(5.5), Inches(1.2))
    set_shape_color(trust_box, BLUE_DARK)
    trust_box.adjustments[0] = 0.15
    trust_txt = slide.shapes.add_textbox(Inches(4), Inches(1.35), Inches(5.5), Inches(1))
    tf = trust_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS TRUST"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Board of progressive leaders"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RgbColor(0xbb, 0xde, 0xfb)
    p2.alignment = PP_ALIGN.CENTER

    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(2.5), Inches(0.7), Inches(0.6))
    set_shape_color(arrow, BLUE_MID)

    llc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(3.2), Inches(5.5), Inches(1))
    set_shape_color(llc_box, BLUE_MID)
    llc_box.adjustments[0] = 0.15
    llc_txt = slide.shapes.add_textbox(Inches(4), Inches(3.4), Inches(5.5), Inches(0.6))
    tf = llc_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS LLC (100% trust-owned)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Left points
    left_points = ["✓ Mission locked to progressives", "✓ Cannot be sold to opponents", "✓ No equity, no hostile exits"]
    for i, point in enumerate(left_points):
        txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.5 + i * 0.6), Inches(3.3), Inches(0.5))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN

    # Right points
    right_points = ["Major funders get board seats", "Transparent governance", "Permanent infrastructure"]
    for i, point in enumerate(right_points):
        txt = slide.shapes.add_textbox(Inches(9.8), Inches(1.5 + i * 0.6), Inches(3.3), Inches(0.5))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE_DARK

    # Catalist precedent
    cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(11.8), Inches(2.6))
    set_shape_color(cat_box, RgbColor(0xf5, 0xf5, 0xf5))
    cat_box.adjustments[0] = 0.05

    cat_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.8), Inches(0.5))
    tf = cat_title.text_frame
    p = tf.paragraphs[0]
    p.text = "THE CATALIST PRECEDENT — Trusted by Democratic campaigns since 2006"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    cat_stats = [("2006", "Founded"), ("$5M+", "Initial funding"), ("Trust LLC", "Structure"), ("Every major\nDem campaign", "Today serves")]
    for i, (num, label) in enumerate(cat_stats):
        x = 1.5 + i * 3
        num_txt = slide.shapes.add_textbox(Inches(x), Inches(5.5), Inches(2.5), Inches(0.7))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        lbl_txt = slide.shapes.add_textbox(Inches(x), Inches(6.2), Inches(2.5), Inches(0.6))
        tf = lbl_txt.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"We're following the Catalist model—trusted by Democratic campaigns since 2006.

Trust owns 100% of the LLC. Mission is locked to progressives. Cannot be sold to opponents.

Major funders get board seats—that's part of the value proposition. You're not buying equity. You're building permanent infrastructure."

TRANSITION: "Let me introduce the team..." """)

    # ========== SLIDE 11: TEAM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Built for Long-Term Impact, Not Quick Exits"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Founder
    founder_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5), Inches(2.8))
    set_shape_color(founder_box, RgbColor(0xf5, 0xf5, 0xf5))
    founder_box.adjustments[0] = 0.08

    photo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.5))
    set_shape_color(photo, BLUE_LIGHT)
    photo_txt = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(1.5), Inches(0.6))
    tf = photo_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "EF"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    f_name = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(2.5), Inches(0.4))
    tf = f_name.text_frame
    p = tf.paragraphs[0]
    p.text = "Ed Forman"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    f_title = slide.shapes.add_textbox(Inches(3), Inches(1.9), Inches(2.5), Inches(0.3))
    tf = f_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Founder & Architect"
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE_MID

    f_details = ["35 years Silicon Valley", "11 startups (early employee)", "Stanford MBA, taught there"]
    for i, detail in enumerate(f_details):
        txt = slide.shapes.add_textbox(Inches(3), Inches(2.4 + i * 0.4), Inches(2.5), Inches(0.3))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + detail
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # Team
    team_title = slide.shapes.add_textbox(Inches(6.3), Inches(1.2), Inches(6.2), Inches(0.4))
    tf = team_title.text_frame
    p = tf.paragraphs[0]
    p.text = "LAUNCH TEAM (Q2 2026)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    roles = [("CEO", "Strategy & fundraising"), ("Tech Lead", "Bentham infrastructure"), ("Research Dir", "Methodology & quality"), ("Ops Manager", "Client delivery"), ("Marketing", "Partnerships & sales")]
    for i, (role, desc) in enumerate(roles):
        x = 6.3 + (i % 3) * 2.1
        y = 1.7 + (i // 3) * 1.3

        role_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.9), Inches(1.1))
        set_shape_color(role_box, RgbColor(0xe3, 0xf2, 0xfd))
        role_box.adjustments[0] = 0.15

        role_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(1.9), Inches(0.4))
        tf = role_txt.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        desc_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.55), Inches(1.9), Inches(0.4))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Board
    board_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.8), Inches(2.9))
    set_shape_color(board_box, RgbColor(0xf5, 0xf5, 0xf5))
    board_box.adjustments[0] = 0.05

    board_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.8), Inches(0.4))
    tf = board_title.text_frame
    p = tf.paragraphs[0]
    p.text = "BOARD OF TRUSTEES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    trustees = [("Progressive\nLeaders", "Mission fidelity"), ("Foundation\nReps", "Major funder seats"), ("Labor\nReps", "Union alignment"), ("Independent\nExpert", "Tech/political")]
    for i, (role, desc) in enumerate(trustees):
        x = 1.2 + i * 3

        t_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.1), Inches(2.5), Inches(1.8))
        set_shape_color(t_box, WHITE, BLUE_MID)
        t_box.adjustments[0] = 0.1
        t_box.line.width = Pt(2)

        t_role = slide.shapes.add_textbox(Inches(x), Inches(5.25), Inches(2.5), Inches(0.8))
        tf = t_role.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        t_desc = slide.shapes.add_textbox(Inches(x), Inches(6.1), Inches(2.5), Inches(0.3))
        tf = t_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"I'm the architect, not the operator. 35 years in Silicon Valley, 11 startups, Stanford MBA.

We're hiring a CEO with deep progressive foundation relationships. Technical lead builds Bentham. Research director ensures quality.

Board includes progressive leaders, foundation reps, labor reps, and independent experts. Major funders get board seats."

TRANSITION: "Here's the growth trajectory..." """)

    # ========== SLIDE 12: GROWTH ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Free Year Builds Reference Base and Proves ROI"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Year cards
    years = [
        ("2026", "FREE YEAR", "200", "campaigns", "$0", "Build reference base", RgbColor(0xe3, 0xf2, 0xfd)),
        ("2027", "LAUNCH", "150", "clients", "$1.2M", "Breakeven", RgbColor(0xc8, 0xe6, 0xc9)),
        ("2028", "SCALE", "466", "clients", "$2.8M", "+$850K net", RgbColor(0xbb, 0xde, 0xfb)),
    ]

    for i, (year, phase, num, label, rev, result, color) in enumerate(years):
        x = 0.8 + i * 4.2

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(3.8), Inches(3.5))
        set_shape_color(card, color)
        card.adjustments[0] = 0.08

        year_txt = slide.shapes.add_textbox(Inches(x), Inches(1.4), Inches(3.8), Inches(0.5))
        tf = year_txt.text_frame
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        phase_txt = slide.shapes.add_textbox(Inches(x), Inches(1.9), Inches(3.8), Inches(0.4))
        tf = phase_txt.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        num_txt = slide.shapes.add_textbox(Inches(x), Inches(2.4), Inches(3.8), Inches(0.9))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        lbl_txt = slide.shapes.add_textbox(Inches(x), Inches(3.3), Inches(3.8), Inches(0.4))
        tf = lbl_txt.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        rev_txt = slide.shapes.add_textbox(Inches(x), Inches(3.8), Inches(3.8), Inches(0.5))
        tf = rev_txt.text_frame
        p = tf.paragraphs[0]
        p.text = rev
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN if i > 0 else GRAY
        p.alignment = PP_ALIGN.CENTER

        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 4), Inches(2.8), Inches(0.4), Inches(0.4))
            set_shape_color(arrow, GRAY)

    # 2026 breakdown
    breakdown_title = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(12), Inches(0.4))
    tf = breakdown_title.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Free Year Breakdown:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    segments = [("15", "Statewide"), ("50", "Congressional"), ("110", "Local"), ("25", "Other orgs")]
    for i, (num, label) in enumerate(segments):
        x = 1 + i * 3
        seg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.5), Inches(2.5), Inches(1.5))
        set_shape_color(seg, RgbColor(0xf5, 0xf5, 0xf5))
        seg.adjustments[0] = 0.15

        seg_num = slide.shapes.add_textbox(Inches(x), Inches(5.6), Inches(2.5), Inches(0.7))
        tf = seg_num.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        seg_lbl = slide.shapes.add_textbox(Inches(x), Inches(6.4), Inches(2.5), Inches(0.3))
        tf = seg_lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"2026 is a free year—we serve 200 campaigns at no cost. 15 statewide, 50 Congressional, 110 local, plus organizations.

Why free? To build reference base and prove ROI. Case studies and testimonials for 2027 sales.

By 2027 we launch paid subscriptions and reach breakeven. By 2028—presidential cycle—466 clients, $2.8M revenue, $850K surplus."

TRANSITION: "Here's what we're asking for..." """)

    # ========== SLIDE 13: THE ASK ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "$1M to Build Progressive AI Infrastructure"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    big_ask = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5), Inches(1.2))
    tf = big_ask.text_frame
    p = tf.paragraphs[0]
    p.text = "$1,000,000"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = BLUE_MID

    range_txt = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(5), Inches(0.6))
    tf = range_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Min: $750K  •  Target: $1M  •  Reach: $1.25M"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    deadline = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(5), Inches(0.5))
    tf = deadline.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 Commitments by March 15, 2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    # Use of funds
    funds_title = slide.shapes.add_textbox(Inches(6), Inches(1.1), Inches(6.5), Inches(0.4))
    tf = funds_title.text_frame
    p = tf.paragraphs[0]
    p.text = "USE OF FUNDS"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    fund_items = [
        ("Personnel (5 FTEs)", "$650K", "65%", BLUE_MID),
        ("Contractors", "$150K", "15%", BLUE_LIGHT),
        ("Infrastructure", "$48K", "5%", GREEN),
        ("Marketing", "$40K", "4%", ORANGE),
        ("Other", "$112K", "11%", GRAY),
    ]

    for i, (name, amount, pct, color) in enumerate(fund_items):
        y = 1.5 + i * 0.6
        bar_width = float(pct.rstrip('%')) / 100 * 5.5
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(y), Inches(bar_width), Inches(0.45))
        set_shape_color(bar, color)
        bar.adjustments[0] = 0.2

        lbl = slide.shapes.add_textbox(Inches(6.1), Inches(y + 0.05), Inches(4), Inches(0.35))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = f"{name}: {amount}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE if bar_width > 1.5 else GRAY

        pct_txt = slide.shapes.add_textbox(Inches(11.7), Inches(y + 0.05), Inches(0.8), Inches(0.35))
        tf = pct_txt.text_frame
        p = tf.paragraphs[0]
        p.text = pct
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # What it enables
    enables_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(11.8), Inches(2.8))
    set_shape_color(enables_box, RgbColor(0xe8, 0xf5, 0xe9))
    enables_box.adjustments[0] = 0.05

    enables_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.8), Inches(0.4))
    tf = enables_title.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT $1M ENABLES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    enables = [("200", "campaigns\nserved free"), ("5", "FTEs\nhired"), ("4", "products\nlaunched"), ("$850K", "surplus\nby 2028")]
    for i, (num, label) in enumerate(enables):
        x = 1.2 + i * 3

        e_num = slide.shapes.add_textbox(Inches(x), Inches(4.9), Inches(2.5), Inches(0.8))
        tf = e_num.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        e_lbl = slide.shapes.add_textbox(Inches(x), Inches(5.8), Inches(2.5), Inches(0.8))
        tf = e_lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"We're asking for $1 million. Min $750K, reach $1.25M. Commitments needed by March 15.

65% goes to personnel—5 FTEs. 15% contractors. Rest is infrastructure, marketing, operations.

What does $1M enable? 200 campaigns served free. 5 FTEs hired. 4 products launched. $850K annual surplus by 2028."

TRANSITION: "Let me tell you why I'm in this room specifically..." """)

    # ========== SLIDE 14: WHY INVESTING IN US ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kyanos Completes Your Portfolio"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Quote
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.1), Inches(11.8), Inches(1.2))
    set_shape_color(quote_box, RgbColor(0xe3, 0xf2, 0xfd))
    quote_box.adjustments[0] = 0.1

    quote_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.8), Inches(0.6))
    tf = quote_txt.text_frame
    p = tf.paragraphs[0]
    p.text = '"Technology is changing politics faster than politics is adapting."'
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    quote_attr = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.8), Inches(0.3))
    tf = quote_attr.text_frame
    p = tf.paragraphs[0]
    p.text = "— Investing In US founding principle"
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE_MID
    p.alignment = PP_ALIGN.CENTER

    # Portfolio fit
    portfolio = [
        ("Vote.org", "Voter registration", "→ We ensure they find accurate info"),
        ("News for Democracy", "Media ecosystem", "→ We ensure AI surfaces accurate news"),
        ("Run for Something", "Candidate pipeline", "→ We help them compete in AI"),
    ]

    for i, (org, does, kyanos) in enumerate(portfolio):
        y = 2.6 + i * 1

        org_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(3), Inches(0.8))
        set_shape_color(org_box, BLUE_LIGHT)
        org_box.adjustments[0] = 0.15

        org_txt = slide.shapes.add_textbox(Inches(1), Inches(y + 0.1), Inches(3), Inches(0.3))
        tf = org_txt.text_frame
        p = tf.paragraphs[0]
        p.text = org
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        does_txt = slide.shapes.add_textbox(Inches(1), Inches(y + 0.45), Inches(3), Inches(0.3))
        tf = does_txt.text_frame
        p = tf.paragraphs[0]
        p.text = does
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        k_txt = slide.shapes.add_textbox(Inches(4.5), Inches(y + 0.2), Inches(7.5), Inches(0.4))
        tf = k_txt.text_frame
        p = tf.paragraphs[0]
        p.text = kyanos
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GREEN

    # CTA
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.4))
    set_shape_color(cta_box, BLUE_DARK)
    cta_box.adjustments[0] = 0.12

    cta_txt = slide.shapes.add_textbox(Inches(0.8), Inches(6), Inches(11.8), Inches(1))
    tf = cta_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Help keep progressives visible"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "when information access determines outcomes."
    p2.font.size = Pt(24)
    p2.font.color.rgb = RgbColor(0xbb, 0xde, 0xfb)
    p2.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Investing In US was founded on this observation: technology changes politics faster than politics adapts.

You've invested in voter registration, candidate pipelines, media ecosystems. Kyanos is the AI layer that makes all of those more effective.

Vote.org registers voters—we ensure they find accurate info when they research. Run for Something builds candidates—we help them compete in AI environments.

Here's my ask: bring this to Reid and Dmitri. If it resonates, I'll give you everything you need to make the case internally.

Help keep progressives visible when information access determines outcomes."

FINAL: "What questions do you have? What would they need to see?" """)

    # Save
    prs.save('/Users/edf/KyanosTech-Plan/docs/presentations/kyanos-improved.pptx')
    print("Created: kyanos-improved.pptx")

if __name__ == "__main__":
    main()
