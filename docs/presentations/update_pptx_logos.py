#!/usr/bin/env python3
"""Add logos to the Kyanos presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as RgbColor

# Open existing presentation
prs = Presentation('/Users/edf/KyanosTech-Plan/docs/presentations/kyanos-visual-pitch.pptx')

# Slide 1 - Add main logo
slide1 = prs.slides[0]
# Remove the placeholder circle (shape index may vary)
# Add the actual logo
logo_path = '/Users/edf/KyanosTech-Plan/docs/presentations/logos/kyanos-main.png'
slide1.shapes.add_picture(logo_path, Inches(5.9), Inches(1.2), height=Inches(2))

# Slide 8 - Add product logos
slide8 = prs.slides[7]  # 0-indexed
logo_size = Inches(0.6)

# Product logo positions (approximate, on product cards)
product_logos = [
    ('/Users/edf/KyanosTech-Plan/docs/presentations/logos/panopticon.png', 1.1, 1.5),
    ('/Users/edf/KyanosTech-Plan/docs/presentations/logos/goodink.png', 7.4, 1.5),
    ('/Users/edf/KyanosTech-Plan/docs/presentations/logos/adwatch.png', 1.1, 4.0),
    ('/Users/edf/KyanosTech-Plan/docs/presentations/logos/tripwire.png', 7.4, 4.0),
]

for logo_path, x, y in product_logos:
    slide8.shapes.add_picture(logo_path, Inches(x), Inches(y), height=logo_size)

# Save updated presentation
prs.save('/Users/edf/KyanosTech-Plan/docs/presentations/kyanos-visual-pitch-with-logos.pptx')
print("Created: kyanos-visual-pitch-with-logos.pptx")
