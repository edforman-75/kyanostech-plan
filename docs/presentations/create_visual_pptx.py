#!/usr/bin/env python3
"""Generate visually engaging Kyanos investor presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

# Color scheme
BLUE_DARK = RgbColor(0x1a, 0x36, 0x5d)    # Navy
BLUE_MID = RgbColor(0x00, 0x7b, 0xc0)      # Kyanos blue
BLUE_LIGHT = RgbColor(0x4a, 0xb8, 0xe8)    # Light blue
ORANGE = RgbColor(0xff, 0x6b, 0x35)        # Accent
GREEN = RgbColor(0x2e, 0xcc, 0x71)         # Success
GRAY = RgbColor(0x6c, 0x75, 0x7d)          # Muted
WHITE = RgbColor(0xff, 0xff, 0xff)
BLACK = RgbColor(0x1a, 0x1a, 0x1a)

def set_shape_color(shape, fill_color, line_color=None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()

def add_big_stat(slide, number, label, x, y, width=3, color=BLUE_MID):
    """Add a big statistic with label below."""
    # Number
    num_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(1.2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(Inches(x), Inches(y + 1.1), Inches(width), Inches(0.6))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_logo_placeholder(slide, name, x, y, size=0.8):
    """Add a placeholder circle for a logo."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    set_shape_color(shape, RgbColor(0xf0, 0xf0, 0xf0), GRAY)

    # Label
    label = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + size + 0.05), Inches(size + 0.6), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_arrow(slide, x1, y1, x2, y2, color=BLUE_MID):
    """Add an arrow shape."""
    # Use a simple line with arrow
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1), Inches(y1),
        Inches(x2 - x1), Inches(0.4)
    )
    set_shape_color(connector, color)

def add_icon_box(slide, icon_text, label, x, y, color=BLUE_MID):
    """Add a colored box with icon text and label."""
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2), Inches(1.5))
    set_shape_color(box, color)
    box.adjustments[0] = 0.1

    # Icon text (emoji or symbol)
    icon = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(2), Inches(0.6))
    tf = icon.text_frame
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl = slide.shapes.add_textbox(Inches(x), Inches(y + 0.8), Inches(2), Inches(0.5))
    tf = lbl.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_notes(slide, notes_text):
    """Add speaker notes to slide."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # ========== SLIDE 1: TITLE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Blue gradient bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.3))
    set_shape_color(bar, BLUE_MID)

    # Logo placeholder
    logo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(1.5), Inches(1.5), Inches(1.5))
    set_shape_color(logo, BLUE_MID)
    logo_text = slide.shapes.add_textbox(Inches(6), Inches(1.8), Inches(1.5), Inches(1))
    tf = logo_text.text_frame
    p = tf.paragraphs[0]
    p.text = "K"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(12.333), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tag = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.6))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = "Keeping Progressives Visible in the AI Era"
    p.font.size = Pt(28)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Presenter info
    info = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "Ed Forman, Founder  •  For Thanasi Dilos, Investing In US  •  January 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Thanasi, thank you for your time. I reached out to Investing In US because you understand something most political operatives don't: the next communication revolution isn't coming—it's already here.

Kyanos is the infrastructure play for progressive politics in the AI era."

TRANSITION: "Let me start with a lesson we should have learned in 2024..." """)

    # ========== SLIDE 2: PODCAST MISTAKE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "We Missed Podcasts. AI Is Next."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Left side - Podcasts (red/warning)
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5))
    set_shape_color(box1, RgbColor(0xff, 0xeb, 0xee))
    box1.adjustments[0] = 0.05

    header1 = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(5.5), Inches(0.6))
    tf = header1.text_frame
    p = tf.paragraphs[0]
    p.text = "2020-2024: PODCASTS"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Joe Rogan stat
    stat1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(5.5), Inches(1))
    tf = stat1.text_frame
    p = tf.paragraphs[0]
    p.text = "14M"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    label1 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.8))
    tf = label1.text_frame
    p = tf.paragraphs[0]
    p.text = "viewers saw Rogan's\nTrump endorsement"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    result1 = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(5.5), Inches(1))
    tf = result1.text_frame
    p = tf.paragraphs[0]
    p.text = "Progressives arrived late"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.4), Inches(3.8), Inches(0.8), Inches(0.5))
    set_shape_color(arrow, BLUE_MID)

    # Right side - AI (green/opportunity)
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.8), Inches(5.5), Inches(4.5))
    set_shape_color(box2, RgbColor(0xe8, 0xf5, 0xe9))
    box2.adjustments[0] = 0.05

    header2 = slide.shapes.add_textbox(Inches(7.2), Inches(2), Inches(5.5), Inches(0.6))
    tf = header2.text_frame
    p = tf.paragraphs[0]
    p.text = "2024-2028: AI"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Logos row
    add_logo_placeholder(slide, "Google", 7.8, 2.7)
    add_logo_placeholder(slide, "ChatGPT", 9.2, 2.7)
    add_logo_placeholder(slide, "Meta AI", 10.6, 2.7)

    result2 = slide.shapes.add_textbox(Inches(7.2), Inches(5), Inches(5.5), Inches(1))
    tf = result2.text_frame
    p = tf.paragraphs[0]
    p.text = "We can lead this time"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"You watched this happen. Joe Rogan's endorsement of Trump reached 14 million people in a single episode. Progressives had no equivalent infrastructure.

The same pattern is emerging with AI. Right now, most campaigns treat AI as a curiosity—something the tech team might mention, not something the campaign manager prioritizes.

That's exactly how we treated podcasts in 2020."

TRANSITION: "But AI isn't just another podcast moment. The data shows it's more consequential..." """)

    # ========== SLIDE 3: AI STATS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The AI Information Revolution"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Big stats row
    add_big_stat(slide, "58%", "of Google searches\nend without a click", 0.5, 1.3, 3)
    add_big_stat(slide, "50%", "of queries show\nAI Overviews", 3.7, 1.3, 3)
    add_big_stat(slide, "100M+", "weekly ChatGPT\nusers in US", 6.9, 1.3, 3)
    add_big_stat(slide, "45-60%", "voter queries by AI\nin 2028", 10.1, 1.3, 3, ORANGE)

    # Platform boxes
    box_y = 4.2
    # Google
    g_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(box_y), Inches(3.5), Inches(2.2))
    set_shape_color(g_box, RgbColor(0xea, 0x43, 0x35))  # Google red
    g_box.adjustments[0] = 0.1
    add_logo_placeholder(slide, "", 2.3, box_y + 0.3, 0.6)
    g_text = slide.shapes.add_textbox(Inches(1), Inches(box_y + 1.1), Inches(3.5), Inches(1))
    tf = g_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Google AI Overviews\nAppears before all links"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # ChatGPT
    c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(box_y), Inches(3.5), Inches(2.2))
    set_shape_color(c_box, RgbColor(0x10, 0xa3, 0x7f))  # OpenAI green
    c_box.adjustments[0] = 0.1
    add_logo_placeholder(slide, "", 6.3, box_y + 0.3, 0.6)
    c_text = slide.shapes.add_textbox(Inches(5), Inches(box_y + 1.1), Inches(3.5), Inches(1))
    tf = c_text.text_frame
    p = tf.paragraphs[0]
    p.text = "ChatGPT\n100M+ active researchers"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # WhatsApp/Meta
    w_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(box_y), Inches(3.5), Inches(2.2))
    set_shape_color(w_box, RgbColor(0x00, 0x84, 0xff))  # Meta blue
    w_box.adjustments[0] = 0.1
    add_logo_placeholder(slide, "", 10.3, box_y + 0.3, 0.6)
    w_text = slide.shapes.add_textbox(Inches(9), Inches(box_y + 1.1), Inches(3.5), Inches(1))
    tf = w_text.text_frame
    p = tf.paragraphs[0]
    p.text = "WhatsApp AI\nCritical for Hispanic voters"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"This isn't about some future scenario. 58% of Google searches already end without a click—voters get their answer from AI Overviews and never visit a campaign website.

By the time voters are making choices in 2028, nearly half of all political information queries will be answered inside AI environments.

The mission of the candidate website has transformed. It's no longer primarily a destination for voters—it's the source material that AI systems draw from."

TRANSITION: "And here's what makes this urgent: AI isn't just reaching voters. It's moving them..." """)

    # ========== SLIDE 4: PERSUASION RESEARCH ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Is 4x More Persuasive Than Television"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Source badge
    source = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(0.4), Inches(3.3), Inches(0.5))
    set_shape_color(source, RgbColor(0xe3, 0xf2, 0xfd))
    source.adjustments[0] = 0.3
    src_text = slide.shapes.add_textbox(Inches(9.5), Inches(0.45), Inches(3.3), Inches(0.4))
    tf = src_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Nature & Science, Dec 2025"
    p.font.size = Pt(12)
    p.font.color.rgb = BLUE_MID
    p.alignment = PP_ALIGN.CENTER

    # Big 4X comparison
    # TV bar
    tv_label = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(2), Inches(0.4))
    tf = tv_label.text_frame
    p = tf.paragraphs[0]
    p.text = "TV Advertising"
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY

    tv_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.3), Inches(2.5), Inches(0.8))
    set_shape_color(tv_bar, GRAY)
    tv_bar.adjustments[0] = 0.2

    tv_val = slide.shapes.add_textbox(Inches(1), Inches(2.4), Inches(2.5), Inches(0.6))
    tf = tv_val.text_frame
    p = tf.paragraphs[0]
    p.text = "1x"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # AI bar
    ai_label = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(2), Inches(0.4))
    tf = ai_label.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Chatbots"
    p.font.size = Pt(18)
    p.font.color.rgb = BLUE_MID

    ai_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4), Inches(10), Inches(0.8))
    set_shape_color(ai_bar, BLUE_MID)
    ai_bar.adjustments[0] = 0.05

    ai_val = slide.shapes.add_textbox(Inches(1), Inches(4.1), Inches(10), Inches(0.6))
    tf = ai_val.text_frame
    p = tf.paragraphs[0]
    p.text = "4x"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Cost comparison boxes
    cost_y = 5.3
    # TV cost
    tv_cost = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(cost_y), Inches(4), Inches(1.5))
    set_shape_color(tv_cost, RgbColor(0xf5, 0xf5, 0xf5))
    tv_cost.adjustments[0] = 0.1

    tv_price = slide.shapes.add_textbox(Inches(1), Inches(cost_y + 0.2), Inches(4), Inches(0.7))
    tf = tv_price.text_frame
    p = tf.paragraphs[0]
    p.text = "$100-600"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    tv_desc = slide.shapes.add_textbox(Inches(1), Inches(cost_y + 0.9), Inches(4), Inches(0.4))
    tf = tv_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "per voter influenced (TV)"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # AI cost
    ai_cost = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(cost_y), Inches(4), Inches(1.5))
    set_shape_color(ai_cost, RgbColor(0xe3, 0xf2, 0xfd))
    ai_cost.adjustments[0] = 0.1

    ai_price = slide.shapes.add_textbox(Inches(5.5), Inches(cost_y + 0.2), Inches(4), Inches(0.7))
    tf = ai_price.text_frame
    p = tf.paragraphs[0]
    p.text = "$1-10"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    ai_desc = slide.shapes.add_textbox(Inches(5.5), Inches(cost_y + 0.9), Inches(4), Inches(0.4))
    tf = ai_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "per voter influenced (AI)"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

    # Quote
    quote = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(0.5))
    tf = quote.text_frame
    p = tf.paragraphs[0]
    p.text = '"A shockingly large effect" — David Rand, MIT'
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = GRAY

    add_notes(slide, """SPEAKER NOTES:

"This is peer-reviewed research from the world's most prestigious scientific journals, published December 2025.

AI chatbots are the most persuasive political communication channel ever measured in controlled experiments. Four times more effective than television advertising.

And the cost? $1-10 per voter influenced, compared to $100-600 for broadcast TV.

This isn't just about keeping pace. It's about finding a lever that actually works."

TRANSITION: "But there's another dimension to this that I know you've been thinking about..." """)

    # ========== SLIDE 5: PLATFORM RISK ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "AI Platforms Face Structural Pressure"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Concentration diagram - companies at top
    company_y = 1.5
    companies = [("Microsoft", 1.5), ("Google", 4.5), ("Amazon", 7.5), ("Meta", 10.5)]
    for name, x in companies:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(company_y), Inches(2), Inches(0.8))
        set_shape_color(box, BLUE_DARK)
        box.adjustments[0] = 0.2
        txt = slide.shapes.add_textbox(Inches(x), Inches(company_y + 0.2), Inches(2), Inches(0.5))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Arrow down
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(2.5), Inches(1), Inches(0.8))
    set_shape_color(arrow, ORANGE)

    # Pressure boxes
    pressure_y = 3.5
    pressures = [
        ("$9B+", "DoD contracts"),
        ("230", "Section 230 risk"),
        ("⚠️", "Political pressure"),
    ]
    for i, (icon, label) in enumerate(pressures):
        x = 2 + i * 3.5
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(pressure_y), Inches(3), Inches(1.5))
        set_shape_color(box, RgbColor(0xff, 0xf3, 0xe0))
        box.adjustments[0] = 0.1

        ico = slide.shapes.add_textbox(Inches(x), Inches(pressure_y + 0.2), Inches(3), Inches(0.6))
        tf = ico.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = ORANGE
        p.alignment = PP_ALIGN.CENTER

        lbl = slide.shapes.add_textbox(Inches(x), Inches(pressure_y + 0.9), Inches(3), Inches(0.4))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Key message
    msg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.5), Inches(11.5), Inches(1.3))
    set_shape_color(msg_box, RgbColor(0xfc, 0xe4, 0xec))
    msg_box.adjustments[0] = 0.1

    msg = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(11.5), Inches(1))
    tf = msg.text_frame
    p = tf.paragraphs[0]
    p.text = "The issue is structural: corporate governance requires managing political risk."
    p.font.size = Pt(22)
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "In the current environment, that risk runs in one direction."
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = ORANGE
    p2.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"You know these companies. You know the pressures they face.

I'm not suggesting bad faith. I'm suggesting that the incentive structures are misaligned. When an administration is willing to punish perceived disloyalty—and when companies have billions in government contracts at stake—caution has a political direction.

That's why external monitoring matters. We need independent visibility into how AI systems represent political topics.

Kyanos provides that watchdog function."

TRANSITION: "So what exactly are we building?" """)

    # ========== SLIDE 6: SOLUTION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kyanos: The Progressive AI Watchdog"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Process flow - 4 steps
    steps = [
        ("🔍", "MONITOR", "See what AI says"),
        ("🔗", "TRACE", "Find root causes"),
        ("📊", "DIAGNOSE", "Rank by impact"),
        ("✓", "FIX", "Actionable steps"),
    ]

    step_width = 2.5
    start_x = 1.2
    step_y = 1.8

    for i, (icon, title_text, desc) in enumerate(steps):
        x = start_x + i * 3

        # Circle with icon
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.6), Inches(step_y), Inches(1.2), Inches(1.2))
        set_shape_color(circle, BLUE_MID)

        icon_txt = slide.shapes.add_textbox(Inches(x + 0.6), Inches(step_y + 0.25), Inches(1.2), Inches(0.8))
        tf = icon_txt.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

        # Title
        step_title = slide.shapes.add_textbox(Inches(x), Inches(step_y + 1.4), Inches(step_width), Inches(0.4))
        tf = step_title.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        # Description
        step_desc = slide.shapes.add_textbox(Inches(x), Inches(step_y + 1.8), Inches(step_width), Inches(0.4))
        tf = step_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 2.3), Inches(step_y + 0.4), Inches(0.5), Inches(0.4))
            set_shape_color(arrow, RgbColor(0xcc, 0xcc, 0xcc))

    # Key differentiators
    diff_y = 4
    diffs = [
        ("For campaigns", "Not generic SEO"),
        ("AI-native", "Not search rankings"),
        ("Root causes", "Not just reports"),
        ("CMS-specific", "Ready to execute"),
    ]

    for i, (main, sub) in enumerate(diffs):
        x = 0.8 + i * 3.2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(diff_y), Inches(2.8), Inches(1.4))
        set_shape_color(box, RgbColor(0xf5, 0xf5, 0xf5))
        box.adjustments[0] = 0.1

        main_txt = slide.shapes.add_textbox(Inches(x), Inches(diff_y + 0.3), Inches(2.8), Inches(0.5))
        tf = main_txt.text_frame
        p = tf.paragraphs[0]
        p.text = main
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        sub_txt = slide.shapes.add_textbox(Inches(x), Inches(diff_y + 0.8), Inches(2.8), Inches(0.4))
        tf = sub_txt.text_frame
        p = tf.paragraphs[0]
        p.text = sub
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Watchdog callout
    watch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.7), Inches(11.8), Inches(1))
    set_shape_color(watch, BLUE_DARK)
    watch.adjustments[0] = 0.15

    watch_txt = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.8), Inches(0.6))
    tf = watch_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "🔔  If AI models shift against progressives, we know first and sound the alarm"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"The name Kyanos comes from the ancient Greek word for blue—the root of 'cyan.' We chose it because we're building something foundational.

Our methodology is simple: Monitor what AI says, Trace where it got that information, Diagnose the root cause, and provide specific Fixes.

We don't just hand over a report. We provide step-by-step instructions tailored to each campaign's CMS—whether that's WordPress, Squarespace, or NationBuilder."

TRANSITION: "Let me show you how this works in practice..." """)

    # ========== SLIDE 7: HOW IT WORKS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Causality Chain in Action"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Flow diagram - top to bottom
    # Query
    q_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.3), Inches(5), Inches(1))
    set_shape_color(q_box, BLUE_LIGHT)
    q_box.adjustments[0] = 0.15
    q_txt = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(0.6))
    tf = q_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "💬  \"What does [Candidate] think about abortion?\""
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Arrow
    a1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.2), Inches(2.4), Inches(0.6), Inches(0.5))
    set_shape_color(a1, GRAY)

    # AI Response (problem)
    r_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3), Inches(5), Inches(1.2))
    set_shape_color(r_box, RgbColor(0xff, 0xcc, 0xbc))
    r_box.adjustments[0] = 0.1
    r_txt = slide.shapes.add_textbox(Inches(1), Inches(3.1), Inches(5), Inches(1))
    tf = r_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "❌ AI Response:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p = tf.add_paragraph()
    p.text = "\"The candidate has not stated a clear position...\""
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY

    # Arrow
    a2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.2), Inches(4.3), Inches(0.6), Inches(0.5))
    set_shape_color(a2, GRAY)

    # Root Cause
    c_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.9), Inches(5), Inches(1.2))
    set_shape_color(c_box, RgbColor(0xff, 0xf9, 0xc4))
    c_box.adjustments[0] = 0.1
    c_txt = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(5), Inches(1))
    tf = c_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "🔍 Root Cause:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p = tf.add_paragraph()
    p.text = "Issues page missing structured data markup"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY

    # Right side - the fix
    fix_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.3), Inches(5.5), Inches(5.5))
    set_shape_color(fix_box, RgbColor(0xe8, 0xf5, 0xe9))
    fix_box.adjustments[0] = 0.05

    fix_title = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = fix_title.text_frame
    p = tf.paragraphs[0]
    p.text = "✓ THE FIX"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Fix steps
    fix_steps = [
        "1. Add structured FAQ to issues page",
        "2. Include abortion position statement",
        "3. Add schema.org markup",
        "4. Submit to Google Search Console",
    ]

    for i, step in enumerate(fix_steps):
        step_txt = slide.shapes.add_textbox(Inches(7.3), Inches(2.2 + i * 0.7), Inches(5), Inches(0.5))
        tf = step_txt.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE_DARK

    # CMS specific
    cms_label = slide.shapes.add_textbox(Inches(7.3), Inches(5.2), Inches(5), Inches(0.4))
    tf = cms_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Instructions tailored for:"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY

    # CMS logos placeholder
    cms_names = ["WordPress", "Squarespace", "NationBuilder"]
    for i, name in enumerate(cms_names):
        cms_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3 + i * 1.7), Inches(5.6), Inches(1.5), Inches(0.6))
        set_shape_color(cms_box, WHITE, GRAY)
        cms_box.adjustments[0] = 0.2
        cms_txt = slide.shapes.add_textbox(Inches(7.3 + i * 1.7), Inches(5.7), Inches(1.5), Inches(0.4))
        tf = cms_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Here's a concrete example. ChatGPT says a candidate hasn't stated a clear position on abortion—but they have, prominently on their issues page.

We trace it: the campaign's issues page exists but lacks structured data markup that AI systems can parse.

Then we tell the campaign exactly what needs to happen. The comms director gets the messaging guidance. The web developer gets the technical specs. We provide templates, code snippets, everything they need.

And we don't walk away. We track implementation and confirm when AI systems update."

TRANSITION: "Beyond Panopticon, we have a complete product suite..." """)

    # ========== SLIDE 8: PRODUCTS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Kyanos Product Suite"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Product cards - 2x2 grid
    products = [
        ("PANOPTICON", "MVP+", "AI Monitoring & Remediation", "Core product: see what AI says, get fixes", BLUE_MID),
        ("GOODINK", "MVP", "Positive Press Discovery", "Surface coverage worth featuring", RgbColor(0x00, 0x96, 0x88)),
        ("ADWATCH", "MVP", "Ad Monitoring", "Track opponent ads across Google & Meta", RgbColor(0x7b, 0x1f, 0xa2)),
        ("TRIPWIRE", "MVP", "Opponent Website Monitoring", "Know when competitors change messaging", RgbColor(0xf5, 0x7c, 0x00)),
    ]

    for i, (name, status, subtitle, desc, color) in enumerate(products):
        row = i // 2
        col = i % 2
        x = 0.8 + col * 6.3
        y = 1.3 + row * 2.5

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(2.2))
        set_shape_color(card, RgbColor(0xfa, 0xfa, 0xfa), color)
        card.adjustments[0] = 0.08
        card.line.width = Pt(3)

        # Color bar on left
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.15), Inches(2.2))
        set_shape_color(bar, color)

        # Status badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 4.5), Inches(y + 0.2), Inches(1), Inches(0.4))
        badge_color = GREEN if status == "MVP+" else RgbColor(0x90, 0xa4, 0xae)
        set_shape_color(badge, badge_color)
        badge.adjustments[0] = 0.4
        badge_txt = slide.shapes.add_textbox(Inches(x + 4.5), Inches(y + 0.25), Inches(1), Inches(0.3))
        tf = badge_txt.text_frame
        p = tf.paragraphs[0]
        p.text = status
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Name
        name_txt = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.3), Inches(4), Inches(0.5))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = color

        # Subtitle
        sub_txt = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.9), Inches(5), Inches(0.4))
        tf = sub_txt.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK

        # Description
        desc_txt = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 1.4), Inches(5), Inches(0.6))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    # Watchtower callout at bottom
    watch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.8))
    set_shape_color(watch, BLUE_DARK)
    watch.adjustments[0] = 0.15

    watch_txt = slide.shapes.add_textbox(Inches(0.8), Inches(6.45), Inches(11.8), Inches(0.5))
    tf = watch_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "🗼 WATCHTOWER: Geographic monitoring network detects regional AI variation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Panopticon is the flagship—it's where campaigns see what AI says about them and get fixes. It's MVP+, meaning we've already built and tested the core functionality.

GoodInk surfaces positive press coverage campaigns could be featuring on their sites.

Adwatch monitors political advertising—what opponents are running, where, and how.

TripWire watches opponent websites. When they update messaging, campaigns know immediately.

All powered by Watchtower, our distributed observation network for geographic monitoring."

TRANSITION: "Now let me explain how we've structured this to be permanent and accountable..." """)

    # ========== SLIDE 9: HOW WE OPERATE ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "How We Operate"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Structure diagram
    # Trust at top
    trust_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(1.2), Inches(5.5), Inches(1.2))
    set_shape_color(trust_box, BLUE_DARK)
    trust_box.adjustments[0] = 0.15
    trust_txt = slide.shapes.add_textbox(Inches(4), Inches(1.3), Inches(5.5), Inches(1))
    tf = trust_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS TRUST"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Board of progressive leaders"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RgbColor(0xbb, 0xde, 0xfb)
    p2.alignment = PP_ALIGN.CENTER

    # Arrow down
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.4), Inches(2.5), Inches(0.7), Inches(0.6))
    set_shape_color(arrow, BLUE_MID)

    # LLC
    llc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(3.2), Inches(5.5), Inches(1))
    set_shape_color(llc_box, BLUE_MID)
    llc_box.adjustments[0] = 0.15
    llc_txt = slide.shapes.add_textbox(Inches(4), Inches(3.4), Inches(5.5), Inches(0.6))
    tf = llc_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS LLC (100% trust-owned)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Key points on sides
    left_points = [
        "✓ Mission locked to progressives",
        "✓ Cannot be sold to opponents",
        "✓ No equity, no exits",
    ]
    for i, point in enumerate(left_points):
        txt = slide.shapes.add_textbox(Inches(0.5), Inches(1.5 + i * 0.6), Inches(3.3), Inches(0.5))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN

    right_points = [
        "Major funders get board seats",
        "Transparent governance",
        "Permanent infrastructure",
    ]
    for i, point in enumerate(right_points):
        txt = slide.shapes.add_textbox(Inches(9.8), Inches(1.5 + i * 0.6), Inches(3.3), Inches(0.5))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(16)
        p.font.color.rgb = BLUE_DARK

    # Catalist precedent box
    cat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(11.8), Inches(2.2))
    set_shape_color(cat_box, RgbColor(0xf5, 0xf5, 0xf5))
    cat_box.adjustments[0] = 0.05

    cat_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.8), Inches(0.5))
    tf = cat_title.text_frame
    p = tf.paragraphs[0]
    p.text = "THE CATALIST PRECEDENT"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    # Catalist stats
    cat_stats = [
        ("2006", "Founded"),
        ("$5M+", "Initial funding"),
        ("Trust LLC", "Structure"),
        ("Every major\nDem campaign", "Today serves"),
    ]
    for i, (num, label) in enumerate(cat_stats):
        x = 1.5 + i * 3
        num_txt = slide.shapes.add_textbox(Inches(x), Inches(5.4), Inches(2.5), Inches(0.6))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        lbl_txt = slide.shapes.add_textbox(Inches(x), Inches(6), Inches(2.5), Inches(0.6))
        tf = lbl_txt.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"I know you've seen a lot of pitch decks. Most promise equity and exits. This is different.

We're not optimizing for an exit. We're building infrastructure that the progressive movement will rely on for decades. The trust structure ensures that.

Catalist proved this model works. They took foundation funding in 2006, built permanent infrastructure, and now serve virtually every major Democratic campaign.

Your investment doesn't come with equity—there is no equity. What you get is a seat at the table for governance, and the knowledge that you helped build something permanent."

TRANSITION: "Here's the market we're going after..." """)

    # ========== SLIDE 10: MARKET ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Path to 466 Clients"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Year timeline
    years = [
        ("2026", "FREE YEAR", "200", "campaigns", "$0", "Build reference base", RgbColor(0xe3, 0xf2, 0xfd)),
        ("2027", "LAUNCH", "150", "clients", "$1.2M", "Breakeven", RgbColor(0xc8, 0xe6, 0xc9)),
        ("2028", "SCALE", "466", "clients", "$2.8M", "+$850K net", RgbColor(0xbb, 0xde, 0xfb)),
    ]

    for i, (year, phase, num, label, rev, result, color) in enumerate(years):
        x = 0.8 + i * 4.2

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(3.8), Inches(3.5))
        set_shape_color(card, color)
        card.adjustments[0] = 0.08

        # Year
        year_txt = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(3.8), Inches(0.5))
        tf = year_txt.text_frame
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        # Phase
        phase_txt = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(3.8), Inches(0.4))
        tf = phase_txt.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        # Big number
        num_txt = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(3.8), Inches(0.9))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_txt = slide.shapes.add_textbox(Inches(x), Inches(3.4), Inches(3.8), Inches(0.4))
        tf = lbl_txt.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        # Revenue
        rev_txt = slide.shapes.add_textbox(Inches(x), Inches(3.9), Inches(3.8), Inches(0.5))
        tf = rev_txt.text_frame
        p = tf.paragraphs[0]
        p.text = rev
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = GREEN if i > 0 else GRAY
        p.alignment = PP_ALIGN.CENTER

        # Result
        res_txt = slide.shapes.add_textbox(Inches(x), Inches(4.4), Inches(3.8), Inches(0.3))
        tf = res_txt.text_frame
        p = tf.paragraphs[0]
        p.text = result
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

        # Arrow between cards
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + 4), Inches(3), Inches(0.4), Inches(0.4))
            set_shape_color(arrow, GRAY)

    # 2026 breakdown
    breakdown_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.1), Inches(12), Inches(0.4))
    tf = breakdown_title.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Free Year Breakdown:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    segments = [("15", "Statewide"), ("50", "Congress"), ("110", "Local"), ("25", "Other")]
    for i, (num, label) in enumerate(segments):
        x = 1 + i * 3
        seg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.5), Inches(2.5), Inches(1.2))
        set_shape_color(seg, RgbColor(0xf5, 0xf5, 0xf5))
        seg.adjustments[0] = 0.15

        seg_num = slide.shapes.add_textbox(Inches(x), Inches(5.6), Inches(2.5), Inches(0.6))
        tf = seg_num.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        seg_lbl = slide.shapes.add_textbox(Inches(x), Inches(6.2), Inches(2.5), Inches(0.3))
        tf = seg_lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"The 2026 free year is strategic. We're investing $330K of the raise to serve 200 campaigns at no cost. That builds our reference base and generates the case studies we need.

By 2027, we launch paid subscriptions and reach breakeven with 150 clients.

By 2028—the presidential cycle—we project 466 clients and $2.8M in revenue, with $850K net.

The go-to-market is deliberately low-friction. Small campaigns get self-service. Larger campaigns come through consultants who already have trusted relationships."

TRANSITION: "You're going to ask about competition..." """)

    # ========== SLIDE 11: COMPETITION ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "First Mover in Progressive AI Monitoring"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Competitive matrix
    # Headers
    headers = ["", "Political\nContext", "AI-\nNative", "Root\nCauses", "Actionable\nFixes", "Progressive\nOnly"]
    for i, h in enumerate(headers):
        x = 0.8 + i * 2
        w = 1.8 if i > 0 else 2.5
        txt = slide.shapes.add_textbox(Inches(x), Inches(1.3), Inches(w), Inches(0.8))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

    # Competitors
    competitors = [
        ("SEMrush / Ahrefs", ["❌", "❌", "❌", "❌", "❌"]),
        ("AI Monitoring Tools", ["❌", "✓", "❌", "❌", "❌"]),
        ("NGP VAN / Catalist", ["✓", "❌", "❌", "❌", "✓"]),
        ("KYANOS", ["✓", "✓", "✓", "✓", "✓"]),
    ]

    for row, (name, checks) in enumerate(competitors):
        y = 2.1 + row * 1.1
        is_kyanos = name == "KYANOS"

        # Row background
        if is_kyanos:
            bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(y - 0.1), Inches(12.2), Inches(1))
            set_shape_color(bg, RgbColor(0xe3, 0xf2, 0xfd))
            bg.adjustments[0] = 0.15

        # Name
        name_txt = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.15), Inches(2.5), Inches(0.6))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(16)
        p.font.bold = is_kyanos
        p.font.color.rgb = BLUE_MID if is_kyanos else GRAY

        # Checks
        for i, check in enumerate(checks):
            x = 3.3 + i * 2
            check_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.1), Inches(1.8), Inches(0.6))
            tf = check_txt.text_frame
            p = tf.paragraphs[0]
            p.text = check
            p.font.size = Pt(24)
            p.font.color.rgb = GREEN if check == "✓" else RgbColor(0xcc, 0xcc, 0xcc)
            p.alignment = PP_ALIGN.CENTER

    # Network effects callout
    net_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.2))
    set_shape_color(net_box, BLUE_DARK)
    net_box.adjustments[0] = 0.1

    net_title = slide.shapes.add_textbox(Inches(0.8), Inches(6), Inches(11.8), Inches(0.4))
    tf = net_title.text_frame
    p = tf.paragraphs[0]
    p.text = "NETWORK EFFECTS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xbb, 0xde, 0xfb)
    p.alignment = PP_ALIGN.CENTER

    net_desc = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.8), Inches(0.5))
    tf = net_desc.text_frame
    p = tf.paragraphs[0]
    p.text = "More campaigns → Better intelligence → Stronger platform accountability evidence"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"You'll notice there's no direct competitor. That's because the category doesn't exist yet.

SEO tools aren't built for AI. AI monitoring tools aren't built for politics. Political tech companies are focused on voter data, not AI presence.

We're the first mover in progressive AI monitoring, and our trust structure means we can't be acquired by someone who would use this against us.

The network effects matter here. Every campaign we serve makes our intelligence better. Every fix we implement teaches us more about how AI systems source information."

TRANSITION: "Let me introduce the team..." """)

    # ========== SLIDE 12: TEAM ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Built for Permanence"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Founder section
    founder_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(5), Inches(2.5))
    set_shape_color(founder_box, RgbColor(0xf5, 0xf5, 0xf5))
    founder_box.adjustments[0] = 0.08

    # Photo placeholder
    photo = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), Inches(1.5), Inches(1.5), Inches(1.5))
    set_shape_color(photo, BLUE_LIGHT)
    photo_txt = slide.shapes.add_textbox(Inches(1.2), Inches(2), Inches(1.5), Inches(0.6))
    tf = photo_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "EF"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Founder info
    f_name = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(2.5), Inches(0.4))
    tf = f_name.text_frame
    p = tf.paragraphs[0]
    p.text = "Ed Forman"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    f_title = slide.shapes.add_textbox(Inches(3), Inches(1.9), Inches(2.5), Inches(0.3))
    tf = f_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Founder & Architect"
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE_MID

    f_details = [
        "35 years Silicon Valley",
        "11 startups (early employee)",
        "Stanford MBA, taught there",
    ]
    for i, detail in enumerate(f_details):
        txt = slide.shapes.add_textbox(Inches(3), Inches(2.4 + i * 0.4), Inches(2.5), Inches(0.3))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + detail
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY

    # Launch team
    team_title = slide.shapes.add_textbox(Inches(6.3), Inches(1.2), Inches(6.2), Inches(0.4))
    tf = team_title.text_frame
    p = tf.paragraphs[0]
    p.text = "LAUNCH TEAM (Q2 2026)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    roles = [
        ("CEO", "Strategy & fundraising"),
        ("Tech Lead", "Bentham infrastructure"),
        ("Research Dir", "Methodology & quality"),
        ("Ops Manager", "Client delivery"),
        ("Marketing", "Partnerships & sales"),
    ]

    for i, (role, desc) in enumerate(roles):
        x = 6.3 + (i % 3) * 2.1
        y = 1.7 + (i // 3) * 1.3

        role_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.9), Inches(1.1))
        set_shape_color(role_box, RgbColor(0xe3, 0xf2, 0xfd))
        role_box.adjustments[0] = 0.15

        role_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(1.9), Inches(0.4))
        tf = role_txt.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        desc_txt = slide.shapes.add_textbox(Inches(x), Inches(y + 0.55), Inches(1.9), Inches(0.4))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(10)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    # Board section
    board_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.2), Inches(11.8), Inches(2.5))
    set_shape_color(board_box, RgbColor(0xf5, 0xf5, 0xf5))
    board_box.adjustments[0] = 0.05

    board_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.8), Inches(0.4))
    tf = board_title.text_frame
    p = tf.paragraphs[0]
    p.text = "BOARD OF TRUSTEES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    trustees = [
        ("Progressive\nLeaders", "Mission fidelity"),
        ("Foundation\nReps", "Major funder seats"),
        ("Labor\nReps", "Union alignment"),
        ("Independent\nExpert", "Tech/political"),
    ]

    for i, (role, desc) in enumerate(trustees):
        x = 1.2 + i * 3

        t_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5), Inches(2.5), Inches(1.4))
        set_shape_color(t_box, WHITE, BLUE_MID)
        t_box.adjustments[0] = 0.1
        t_box.line.width = Pt(2)

        t_role = slide.shapes.add_textbox(Inches(x), Inches(5.15), Inches(2.5), Inches(0.7))
        tf = t_role.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_DARK
        p.alignment = PP_ALIGN.CENTER

        t_desc = slide.shapes.add_textbox(Inches(x), Inches(5.85), Inches(2.5), Inches(0.3))
        tf = t_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"I'm the architect, not the operator. My role is to assemble the team, frame the problem, and design a structure that outlasts any individual.

We'll hire a CEO with deep relationships in the progressive funding community. The technical lead will build Bentham and the product suite. Research director ensures methodology rigor.

The board of trustees governs the trust and ensures mission fidelity. Major funders get board seats—that's part of the value proposition for foundation investors."

TRANSITION: "Here's what we're asking for..." """)

    # ========== SLIDE 13: THE ASK ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "$1M to Build Progressive AI Infrastructure"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Big number
    big_ask = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5), Inches(1.5))
    tf = big_ask.text_frame
    p = tf.paragraphs[0]
    p.text = "$1,000,000"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = BLUE_MID

    # Range
    range_txt = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(5), Inches(0.8))
    tf = range_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Min: $750K  •  Target: $1M  •  Reach: $1.25M"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY

    deadline = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(5), Inches(0.5))
    tf = deadline.text_frame
    p = tf.paragraphs[0]
    p.text = "📅 Commitments by March 15, 2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    # Use of funds - pie chart style boxes
    funds_title = slide.shapes.add_textbox(Inches(6), Inches(1.2), Inches(6.5), Inches(0.4))
    tf = funds_title.text_frame
    p = tf.paragraphs[0]
    p.text = "USE OF FUNDS"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    fund_items = [
        ("Personnel", "$650K", "65%", BLUE_MID),
        ("Contractors", "$150K", "15%", BLUE_LIGHT),
        ("Infrastructure", "$48K", "5%", GREEN),
        ("Marketing", "$40K", "4%", ORANGE),
        ("Other", "$112K", "11%", GRAY),
    ]

    for i, (name, amount, pct, color) in enumerate(fund_items):
        y = 1.7 + i * 0.7

        # Color bar
        bar_width = float(pct.rstrip('%')) / 100 * 5.5
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(y), Inches(bar_width), Inches(0.5))
        set_shape_color(bar, color)
        bar.adjustments[0] = 0.2

        # Label
        lbl = slide.shapes.add_textbox(Inches(6.1), Inches(y + 0.05), Inches(3), Inches(0.4))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = f"{name}: {amount}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE if bar_width > 1 else GRAY

        # Pct
        pct_txt = slide.shapes.add_textbox(Inches(11.5), Inches(y + 0.05), Inches(1), Inches(0.4))
        tf = pct_txt.text_frame
        p = tf.paragraphs[0]
        p.text = pct
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY

    # What it enables
    enables_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(11.8), Inches(2.5))
    set_shape_color(enables_box, RgbColor(0xe8, 0xf5, 0xe9))
    enables_box.adjustments[0] = 0.05

    enables_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.8), Inches(0.4))
    tf = enables_title.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT $1M ENABLES"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    enables = [
        ("200", "campaigns\nserved free"),
        ("5", "FTEs\nhired"),
        ("4", "products\nlaunched"),
        ("$850K", "surplus\nby 2028"),
    ]

    for i, (num, label) in enumerate(enables):
        x = 1.2 + i * 3

        e_num = slide.shapes.add_textbox(Inches(x), Inches(5), Inches(2.5), Inches(0.8))
        tf = e_num.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = BLUE_MID
        p.alignment = PP_ALIGN.CENTER

        e_lbl = slide.shapes.add_textbox(Inches(x), Inches(5.8), Inches(2.5), Inches(0.8))
        tf = e_lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"We're asking for $1 million. That funds a full year of building and a free access program for 200 campaigns in 2026.

Why can $1 million do what would have taken $5 million a few years ago? AI-assisted development has fundamentally changed the economics. We practice what we preach—Kyanos is built with the same AI capabilities we help campaigns leverage.

The timeline is tight by design. We need commitments by March 15 to be operational by June and serving campaigns before the 2026 cycle heats up."

TRANSITION: "Let me tell you why I'm in this room specifically..." """)

    # ========== SLIDE 14: WHY INVESTING IN US ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "This Is What You Were Built For"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK

    # Quote box
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.2), Inches(11.8), Inches(1.3))
    set_shape_color(quote_box, RgbColor(0xe3, 0xf2, 0xfd))
    quote_box.adjustments[0] = 0.1

    quote_txt = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(0.6))
    tf = quote_txt.text_frame
    p = tf.paragraphs[0]
    p.text = '"Technology is changing politics faster than politics is adapting to technology."'
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    quote_attr = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.8), Inches(0.4))
    tf = quote_attr.text_frame
    p = tf.paragraphs[0]
    p.text = "— Investing In US founding principle"
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE_MID
    p.alignment = PP_ALIGN.CENTER

    # Portfolio fit
    fit_title = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.8), Inches(0.4))
    tf = fit_title.text_frame
    p = tf.paragraphs[0]
    p.text = "KYANOS COMPLETES YOUR PORTFOLIO"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE_DARK
    p.alignment = PP_ALIGN.CENTER

    # Portfolio items
    portfolio = [
        ("Vote.org", "Voter registration", "We ensure they find accurate info"),
        ("News for Democracy", "Media ecosystem", "We ensure AI surfaces accurate news"),
        ("Run for Something", "Candidate pipeline", "We help them compete in AI"),
    ]

    for i, (org, does, kyanos) in enumerate(portfolio):
        y = 3.3 + i * 1

        # Org box
        org_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(3), Inches(0.8))
        set_shape_color(org_box, BLUE_LIGHT)
        org_box.adjustments[0] = 0.15

        org_txt = slide.shapes.add_textbox(Inches(1), Inches(y + 0.1), Inches(3), Inches(0.3))
        tf = org_txt.text_frame
        p = tf.paragraphs[0]
        p.text = org
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        does_txt = slide.shapes.add_textbox(Inches(1), Inches(y + 0.45), Inches(3), Inches(0.3))
        tf = does_txt.text_frame
        p = tf.paragraphs[0]
        p.text = does
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Arrow
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(y + 0.25), Inches(0.5), Inches(0.3))
        set_shape_color(arrow, GRAY)

        # Kyanos addition
        k_txt = slide.shapes.add_textbox(Inches(5), Inches(y + 0.2), Inches(7), Inches(0.4))
        tf = k_txt.text_frame
        p = tf.paragraphs[0]
        p.text = "→ " + kyanos
        p.font.size = Pt(14)
        p.font.color.rgb = GREEN

    # Final CTA
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.2), Inches(11.8), Inches(1))
    set_shape_color(cta_box, BLUE_DARK)
    cta_box.adjustments[0] = 0.15

    cta_txt = slide.shapes.add_textbox(Inches(0.8), Inches(6.4), Inches(11.8), Inches(0.6))
    tf = cta_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Help us keep progressives visible in the AI era."
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_notes(slide, """SPEAKER NOTES:

"Thanasi, Investing In US was founded on a simple observation: technology is changing politics faster than politics is adapting.

That's exactly what we're seeing with AI. Political operatives are treating chatbots like a curiosity while conservatives are already optimizing for them.

You've already invested in voter registration, candidate pipelines, and media ecosystems. Kyanos is the missing layer—ensuring that when voters use AI to research candidates, they find accurate information.

Here's what I know about Investing In US: you're not looking for financial returns. You're looking for democratic returns. You want to build infrastructure that outlasts any single campaign or election cycle.

That's what Kyanos is. Trust-owned. Mission-locked. Built to serve progressives permanently.

Here's my ask: I'd like you to bring this to Reid and Dmitri. If this resonates with you—if you see how it completes your portfolio—I want to give you everything you need to make the case internally.

What questions do you have? What would they need to see?" """)

    # Save
    prs.save('/Users/edf/KyanosTech-Plan/docs/presentations/kyanos-visual-pitch.pptx')
    print("Created: kyanos-visual-pitch.pptx")

if __name__ == "__main__":
    main()
